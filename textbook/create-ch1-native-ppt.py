#!/usr/bin/env python3
"""
Chapter 1 PPT with NATIVE graphics (editable shapes, not images)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# Colors
DARK_SLATE = RGBColor(0x44, 0x54, 0x6A)
SAP_ORANGE = RGBColor(0xED, 0x7D, 0x31)
SAP_BLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHT_GRAY = RGBColor(0xE7, 0xE6, 0xE6)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
GREEN = RGBColor(0x70, 0xAD, 0x47)
DARK_GREEN = RGBColor(0x1D, 0x6F, 0x5C)
DARK_RED = RGBColor(0x8B, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFC, 0xF8, 0xE8)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def add_text(shape, text, size=18, color=DARK_SLATE, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return tf

def add_para(tf, text, size=14, color=DARK_SLATE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    return p

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ==================== SLIDE 1: Title ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
add_text(t, "THE FUTURE OF PRODUCT MARKETING", 13, WHITE, True)

hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(13.333), Inches(4.8))
set_fill(hero, SAP_BLUE); no_line(hero)

tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(2.2), Inches(0.45))
set_fill(tag, SAP_ORANGE); no_line(tag)
tt = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(2.2), Inches(0.4))
add_text(tt, "PART I · CHAPTER 1", 11, WHITE, True, PP_ALIGN.CENTER)

title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(10), Inches(1.8))
add_text(title, "The Old Playbook\nIs Dead", 48, WHITE, True)

sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(10), Inches(0.8))
add_text(sub, "Pragmatic Remix: All 37 Boxes (Reframed)", 16, WHITE)

foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
add_text(foot, "Chris O'Hara & Dan Yu  ·  futureofpmm.com  ·  March 2026", 11, DARK_SLATE)

# ==================== SLIDE 2: Core Thesis ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "CORE THESIS", 20, WHITE, True)

box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3))
set_fill(box, LIGHT_GRAY); no_line(box)
txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.5))
add_text(txt, "The agentic era didn't accelerate the old B2B playbook — it ended it.\n\nProduct lifecycles now run in continuous deployment, not quarterly releases.\nBuyers consume content through AI intermediaries, not your website.\nGEO has replaced SEO. The funnel is no longer a funnel.\n\nThe PMM who survives this transition is not the one who learned to prompt faster —\nit's the one who understood what changed and why.", 16, DARK_SLATE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 3: NATIVE Buyer Journey Compression ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE BUYER JOURNEY COMPRESSION", 20, WHITE, True)

# Old Reality label
old_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4), Inches(0.4))
add_text(old_label, "OLD REALITY (HUMAN-LED)", 11, DARK_SLATE, True)

time1 = slide.shapes.add_textbox(Inches(11.5), Inches(1.1), Inches(1.5), Inches(0.4))
add_text(time1, "~6 weeks", 11, DARK_RED, True)

# Old journey boxes
old_steps = ["Problem\nawareness", "Google\nsearch", "Content\nconsumption", "Vendor\nshortlist", "Demo\nrequests", "Internal\nalignment", "Decision"]
for i, step in enumerate(old_steps):
    x = Inches(0.5 + i*1.75)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(1.6), Inches(0.9))
    set_fill(box, LIGHT_BLUE_BG); no_line(box)
    txt = slide.shapes.add_textbox(x + Inches(0.05), Inches(1.7), Inches(1.5), Inches(0.75))
    add_text(txt, step, 9, DARK_SLATE, False, PP_ALIGN.CENTER)

# Compression arrow/label
comp = slide.shapes.add_textbox(Inches(4), Inches(2.7), Inches(6), Inches(0.4))
add_text(comp, "COMPRESSED BY AI AGENTS → 90 seconds", 11, DARK_RED, True, PP_ALIGN.CENTER)

# New Reality label
new_label = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(4), Inches(0.4))
add_text(new_label, "NEW REALITY (AGENT-MEDIATED)", 11, DARK_GREEN, True)

time2 = slide.shapes.add_textbox(Inches(11.5), Inches(3.3), Inches(1.5), Inches(0.4))
add_text(time2, "~Days", 11, DARK_GREEN, True)

# New journey boxes
new_steps = ["Problem +\nconstraints\nto agent", "Agent\nsynthesizes\nmarket", "Structured\ncomparison\ndelivered", "Human\nreviews\noutput", "Targeted\ndemos\n(if any)", "Decision"]
for i, step in enumerate(new_steps):
    x = Inches(0.5 + i*2.1)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.8), Inches(1.95), Inches(1))
    set_fill(box, LIGHT_GREEN); no_line(box)
    txt = slide.shapes.add_textbox(x + Inches(0.05), Inches(3.9), Inches(1.85), Inches(0.9))
    add_text(txt, step, 9, DARK_SLATE, False, PP_ALIGN.CENTER)

# Bottom callout
callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1))
set_fill(callout, LIGHT_YELLOW); no_line(callout)
ct = slide.shapes.add_textbox(Inches(1.7), Inches(5.5), Inches(9.9), Inches(0.7))
add_text(ct, "Steps 2–4 of the old journey — where most PMM content lived —\nnow happen inside an AI's context window in seconds.", 12, DARK_SLATE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 4: NATIVE Visibility Gap ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE VISIBILITY GAP", 20, WHITE, True)

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12), Inches(0.4))
add_text(subtitle, "What AI agents surface vs. what most PMM teams produce", 12, DARK_SLATE, False, PP_ALIGN.CENTER)

# Left column header - What Agents CAN Process
left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.9), Inches(0.6))
set_fill(left_header, DARK_GREEN); no_line(left_header)
lh_txt = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(5.9), Inches(0.5))
add_text(lh_txt, "WHAT AGENTS CAN PROCESS", 12, WHITE, True, PP_ALIGN.CENTER)

# Right column header - What Agents SKIP
right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.9), Inches(0.6))
set_fill(right_header, DARK_RED); no_line(right_header)
rh_txt = slide.shapes.add_textbox(Inches(6.9), Inches(1.55), Inches(5.9), Inches(0.5))
add_text(rh_txt, "WHAT AGENTS SKIP", 12, WHITE, True, PP_ALIGN.CENTER)

# Left items (what agents can process)
left_items = [
    "Structured comparison data (schema.org)",
    "Third-party review rankings (G2, Capterra)",
    "Technical documentation & APIs",
    "Machine-readable pricing structures",
    "News with verifiable citations"
]

for i, item in enumerate(left_items):
    y = Inches(2.3 + i*0.7)
    check = slide.shapes.add_textbox(Inches(0.6), y, Inches(0.4), Inches(0.5))
    add_text(check, "✓", 14, DARK_GREEN, True)
    txt = slide.shapes.add_textbox(Inches(1), y, Inches(5.3), Inches(0.5))
    add_text(txt, item, 11, DARK_SLATE)

# Right items (what agents skip)
right_items = [
    "Narrative product pages",
    "Video testimonials & demos",
    "Gated PDF case studies & ebooks",
    "JavaScript-rendered content",
    "Email nurture sequences"
]

for i, item in enumerate(right_items):
    y = Inches(2.3 + i*0.7)
    x_mark = slide.shapes.add_textbox(Inches(7), y, Inches(0.4), Inches(0.5))
    add_text(x_mark, "✗", 14, DARK_RED, True)
    txt = slide.shapes.add_textbox(Inches(7.4), y, Inches(5.3), Inches(0.5))
    add_text(txt, item, 11, DARK_SLATE)

# Bottom callout
callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6), Inches(10.3), Inches(0.8))
set_fill(callout, LIGHT_GREEN); no_line(callout)
ct = slide.shapes.add_textbox(Inches(1.7), Inches(6.15), Inches(9.9), Inches(0.6))
add_text(ct, "The Visibility Gap is not a technology problem — it is a prioritization problem.", 13, DARK_GREEN, True, PP_ALIGN.CENTER)

# ==================== SLIDE 5: NATIVE SEO to GEO ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "FROM SEO TO GEO", 20, WHITE, True)

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12), Inches(0.4))
add_text(subtitle, "Generative Engine Optimization", 12, DARK_SLATE, False, PP_ALIGN.CENTER)

# Table header
header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.6))
set_fill(header, DARK_SLATE); no_line(header)

h1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(3.5), Inches(0.5))
add_text(h1, "", 12, WHITE, True)
h2 = slide.shapes.add_textbox(Inches(4.3), Inches(1.55), Inches(3.8), Inches(0.5))
add_text(h2, "SEO", 14, WHITE, True, PP_ALIGN.CENTER)
arrow = slide.shapes.add_textbox(Inches(8.1), Inches(1.55), Inches(0.5), Inches(0.5))
add_text(arrow, "→", 14, GOLD, True, PP_ALIGN.CENTER)
h3 = slide.shapes.add_textbox(Inches(8.6), Inches(1.55), Inches(3.9), Inches(0.5))
add_text(h3, "GEO", 14, DARK_GREEN, True, PP_ALIGN.CENTER)

# Table rows
rows = [
    ("Optimize for", "Ranking on a list", "Citation in a synthesis"),
    ("Compete on", "Keywords", "Entity authority"),
    ("Build", "Backlinks", "Structured proof points"),
    ("What matters", "Title tag", "Schema markup"),
    ("Volume effect", "More content helps", "Better structure helps"),
    ("Buyer sees", "Blue links", "Synthesized answer"),
]

for i, (label, seo, geo) in enumerate(rows):
    y = Inches(2.2 + i*0.65)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    
    row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(0.6))
    set_fill(row_bg, bg_color); no_line(row_bg)
    
    l = slide.shapes.add_textbox(Inches(1), y + Inches(0.1), Inches(3.2), Inches(0.4))
    add_text(l, label, 11, DARK_SLATE, True)
    
    s = slide.shapes.add_textbox(Inches(4.3), y + Inches(0.1), Inches(3.8), Inches(0.4))
    add_text(s, seo, 11, DARK_SLATE, False, PP_ALIGN.CENTER)
    
    g = slide.shapes.add_textbox(Inches(8.6), y + Inches(0.1), Inches(3.9), Inches(0.4))
    add_text(g, geo, 11, DARK_GREEN, True, PP_ALIGN.CENTER)

# Action box
act = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.3), Inches(1))
set_fill(act, GREEN); no_line(act)
at = slide.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(12), Inches(0.35))
add_text(at, "The Audit — Do This Week", 11, WHITE, True)
ad = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(11.9), Inches(0.45))
add_text(ad, "Open Claude, ChatGPT, or Perplexity. Ask it to evaluate vendors in your category. See if your company appears.", 10, WHITE)

# ==================== SLIDE 6: NATIVE Awareness Stack ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE NEW AWARENESS STACK", 20, WHITE, True)

# Layer 1: GEO Layer (green)
layer1_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(0.15), Inches(1.4))
set_fill(layer1_bar, DARK_GREEN); no_line(layer1_bar)
layer1_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.3), Inches(12.2), Inches(1.4))
set_fill(layer1_bg, LIGHT_GREEN); no_line(layer1_bg)

l1_num = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(0.8), Inches(0.6))
add_text(l1_num, "01", 28, DARK_GREEN, True)
l1_title = slide.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(2.5), Inches(0.5))
add_text(l1_title, "GEO LAYER", 14, DARK_GREEN, True)
l1_sub = slide.shapes.add_textbox(Inches(1.8), Inches(1.95), Inches(2.5), Inches(0.4))
add_text(l1_sub, "Technical discoverability", 10, DARK_GREEN)
l1_items = slide.shapes.add_textbox(Inches(5), Inches(1.7), Inches(7.5), Inches(0.6))
add_text(l1_items, "Schema.org markup • Structured data • Machine-readable pricing • FAQ pages mapped to buyer queries", 10, DARK_SLATE)

# Layer 2: Data Layer (blue)
layer2_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.9), Inches(0.15), Inches(1.4))
set_fill(layer2_bar, DARK_SLATE); no_line(layer2_bar)
layer2_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(2.9), Inches(12.2), Inches(1.4))
set_fill(layer2_bg, LIGHT_BLUE_BG); no_line(layer2_bg)

l2_num = slide.shapes.add_textbox(Inches(0.9), Inches(3.1), Inches(0.8), Inches(0.6))
add_text(l2_num, "02", 28, DARK_SLATE, True)
l2_title = slide.shapes.add_textbox(Inches(1.8), Inches(3.1), Inches(2.5), Inches(0.5))
add_text(l2_title, "DATA LAYER", 14, DARK_SLATE, True)
l2_sub = slide.shapes.add_textbox(Inches(1.8), Inches(3.55), Inches(2.5), Inches(0.4))
add_text(l2_sub, "Third-party validation", 10, DARK_SLATE)
l2_items = slide.shapes.add_textbox(Inches(5), Inches(3.3), Inches(7.5), Inches(0.6))
add_text(l2_items, "G2 • Gartner Peer Insights • Capterra • Treat review management with analyst relations rigor", 10, DARK_SLATE)

# Layer 3: Narrative Layer (gold)
layer3_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(0.15), Inches(1.4))
set_fill(layer3_bar, GOLD); no_line(layer3_bar)
layer3_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(4.5), Inches(12.2), Inches(1.4))
set_fill(layer3_bg, LIGHT_YELLOW); no_line(layer3_bg)

l3_num = slide.shapes.add_textbox(Inches(0.9), Inches(4.7), Inches(0.8), Inches(0.6))
add_text(l3_num, "03", 28, GOLD, True)
l3_title = slide.shapes.add_textbox(Inches(1.8), Inches(4.7), Inches(2.5), Inches(0.5))
add_text(l3_title, "NARRATIVE LAYER", 14, GOLD, True)
l3_sub = slide.shapes.add_textbox(Inches(1.8), Inches(5.15), Inches(2.5), Inches(0.4))
add_text(l3_sub, "Irreplaceable human work", 10, GOLD)
l3_items = slide.shapes.add_textbox(Inches(5), Inches(4.9), Inches(7.5), Inches(0.6))
add_text(l3_items, "Category creation • Competitive differentiation • Executive thought leadership • Event keynotes", 10, DARK_SLATE)

# ==================== SLIDE 7: NATIVE Tiered Launch Model ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE TIERED LAUNCH MODEL", 20, WHITE, True)

# Tier 1: Strategic moment
t1_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.08))
set_fill(t1_bar, DARK_SLATE); no_line(t1_bar)
t1_label = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(1.5), Inches(0.4))
add_text(t1_label, "TIER 1", 10, DARK_SLATE)
t1_title = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(3), Inches(0.5))
add_text(t1_title, "Strategic moment", 16, DARK_SLATE, True)
t1_freq = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(2), Inches(0.6))
add_text(t1_freq, "~2x/year", 14, DARK_SLATE, True, PP_ALIGN.CENTER)
t1_desc = slide.shapes.add_textbox(Inches(7), Inches(1.7), Inches(5.5), Inches(0.8))
add_text(t1_desc, "Full positioning refresh, analyst briefing,\npress, keynote, demand gen, sales enablement", 10, DARK_SLATE)

# Tier 2: Feature story
t2_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.08))
set_fill(t2_bar, GREEN); no_line(t2_bar)
t2_label = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(1.5), Inches(0.4))
add_text(t2_label, "TIER 2", 10, DARK_SLATE)
t2_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.6), Inches(3), Inches(0.5))
add_text(t2_title, "Feature story", 16, DARK_SLATE, True)
t2_freq = slide.shapes.add_textbox(Inches(4.5), Inches(3.4), Inches(2), Inches(0.6))
add_text(t2_freq, "Monthly", 14, GREEN, True, PP_ALIGN.CENTER)
t2_desc = slide.shapes.add_textbox(Inches(7), Inches(3.3), Inches(5.5), Inches(0.8))
add_text(t2_desc, "Blog post, demo update, sales talking\npoints, social push, comparison page update", 10, DARK_SLATE)

# Tier 3: Continuous stream
t3_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.08))
set_fill(t3_bar, GOLD); no_line(t3_bar)
t3_label = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(1.5), Inches(0.4))
add_text(t3_label, "TIER 3", 10, DARK_SLATE)
t3_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(3), Inches(0.5))
add_text(t3_title, "Continuous stream", 16, DARK_SLATE, True)
t3_freq = slide.shapes.add_textbox(Inches(4.5), Inches(5), Inches(2), Inches(0.6))
add_text(t3_freq, "Always", 14, GOLD, True, PP_ALIGN.CENTER)
t3_desc = slide.shapes.add_textbox(Inches(7), Inches(4.9), Inches(5.5), Inches(0.8))
add_text(t3_desc, "Structured changelog (human + agent\nreadable), competitive cadence proof", 10, DARK_SLATE)

# ==================== SLIDE 8: Three Clusters ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE PRAGMATIC FRAMEWORK MEETS THE MACHINE", 18, WHITE, True)

sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
add_text(sub, "37 boxes. Three clusters. Which activities survive automation?", 14, DARK_SLATE)

clusters = [
    ("CLUSTER 1", "Automated", "Market sizing, derivative content,\nevent logistics, sales tools.\n\nAgents do the majority today.", LIGHT_GRAY, DARK_SLATE),
    ("CLUSTER 2", "Augmented", "Positioning, win/loss analysis,\nbuyer personas, pricing.\n\nHuman essential, agent accelerates.", SAP_BLUE, WHITE),
    ("CLUSTER 3", "Elevated", "Stakeholder management,\ncustomer empathy, thought leadership.\n\nIrreducibly human.", SAP_ORANGE, WHITE),
]

for i, (num, title, desc, bg, txt_color) in enumerate(clusters):
    x = Inches(0.5 + i*4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(4), Inches(4))
    set_fill(box, bg); no_line(box)
    
    n = slide.shapes.add_textbox(x + Inches(0.2), Inches(2), Inches(3.6), Inches(0.4))
    add_text(n, num, 10, txt_color, True)
    
    tt = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.5))
    add_text(tt, title, 18, txt_color, True)
    
    d = slide.shapes.add_textbox(x + Inches(0.2), Inches(3), Inches(3.6), Inches(2.5))
    add_text(d, desc, 12, txt_color)

ins = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6), Inches(12.3), Inches(1))
set_fill(ins, LIGHT_GRAY); no_line(ins)
it = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7))
add_text(it, "If your value is producing artifacts, the automation wave is not your friend.\nThe question: are you directing the agents or being replaced by them?", 12, DARK_SLATE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 9: 10x Thesis ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, SAP_ORANGE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE 10x THESIS", 20, WHITE, True)

q = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2))
set_fill(q, LIGHT_GRAY); no_line(q)
qt = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.8))
add_text(qt, "A PMM who automates 20 hours a week of Cluster One work and redirects that time into deeper competitive analysis, more thoughtful positioning, and genuine thought leadership will outperform one still manually updating battlecards.\n\nNot by 20%. By a factor that justifies the title of this curriculum.", 15, DARK_SLATE, False, PP_ALIGN.CENTER)

mt = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12), Inches(0.5))
add_text(mt, "What \"10x yourself\" actually means:", 14, DARK_SLATE, True)

meanings = [
    ("✗  NOT working ten times harder", RGBColor(0x8B, 0x00, 0x00)),
    ("✗  NOT working ten times faster", RGBColor(0x8B, 0x00, 0x00)),
    ("✓  Operating at a fundamentally different level of strategic impact", GREEN),
]
for i, (txt, color) in enumerate(meanings):
    row = slide.shapes.add_textbox(Inches(0.5), Inches(4.1 + i*0.5), Inches(12), Inches(0.5))
    add_text(row, txt, 14, color, i==2)

bottom = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1))
add_text(bottom, "You've offloaded the mechanical work to machines built for it —\nand you're finally free to do the work that only a human can do.", 14, DARK_SLATE, True, PP_ALIGN.CENTER)

# ==================== SLIDE 10: CMO Perspective ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE); no_line(h)
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "CMO PERSPECTIVE", 20, WHITE, True)

q = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.6))
set_fill(q, LIGHT_GRAY); no_line(q)
qt = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.4))
add_text(qt, "The conversation about AI inside marketing departments is almost entirely wrong.\nIt's focused on efficiency — that's a CFO question.\n\nThe CMO question: how do we produce fundamentally different outputs?", 14, DARK_SLATE, False, PP_ALIGN.CENTER)

kt = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(0.5))
add_text(kt, "KEY TAKEAWAYS", 14, DARK_SLATE, True)

takeaways = [
    "The agentic shift is a structural transformation, not a toolkit upgrade",
    "Buyer behavior has already changed — AI agents compress research from weeks to hours",
    "Most PMM content is invisible to agents (the Visibility Gap)",
    "The Pragmatic 37 cluster into: Automated, Augmented, Elevated",
    "The PMM who frees Cluster One hours to invest in Cluster Three operates at a different level",
]
for i, txt in enumerate(takeaways):
    color = SAP_BLUE if i < 3 else SAP_ORANGE
    row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.5 + i*0.55), Inches(12.3), Inches(0.5))
    set_fill(row, color); no_line(row)
    rt = slide.shapes.add_textbox(Inches(0.7), Inches(3.6 + i*0.55), Inches(11.9), Inches(0.4))
    add_text(rt, txt, 11, WHITE)

foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), Inches(13.333), Inches(0.7))
set_fill(foot, LIGHT_GRAY); no_line(foot)
ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4))
add_text(ft, "futureofpmm.com  ·  Chris O'Hara & Dan Yu  ·  The Future of Product Marketing", 10, DARK_SLATE, False, PP_ALIGN.CENTER)

# Save
out = "/root/.openclaw/workspace/futureofpmm/dist/curriculum/chapter-01/chapter-01-sap.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  10 slides with ALL NATIVE GRAPHICS (editable shapes)")
