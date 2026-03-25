#!/usr/bin/env python3
"""
Create SAP-styled PowerPoint for Chapter 1
SAP Template Colors:
- Dark slate blue: #44546A (headers, text)
- SAP Orange: #ED7D31 (accents)
- Blue: #4472C4 (secondary)
- Light gray: #E7E6E6 (backgrounds)
- Gold: #FFC000 (highlights)
- Light blue: #5B9BD5 (gradients)
- Green: #70AD47 (success/action)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# SAP Colors
DARK_SLATE = RGBColor(0x44, 0x54, 0x6A)
SAP_ORANGE = RGBColor(0xED, 0x7D, 0x31)
SAP_BLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHT_GRAY = RGBColor(0xE7, 0xE6, 0xE6)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
GREEN = RGBColor(0x70, 0xAD, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_shape_fill(shape, color):
    """Set solid fill color for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text_frame(shape, text, font_size=18, font_color=DARK_SLATE, bold=False, alignment=PP_ALIGN.LEFT):
    """Configure text in a shape"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = alignment
    return tf

def add_paragraph(text_frame, text, font_size=14, font_color=DARK_SLATE, bold=False, bullet=False):
    """Add a paragraph to existing text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Arial"
    if bullet:
        p.level = 0
    return p

# Create presentation (16:9 widescreen)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Get blank layout
blank_layout = prs.slide_layouts[6]

# ========== SLIDE 1: Title Slide ==========
slide = prs.slides.add_slide(blank_layout)

# Dark header bar
header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()

# Title text
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
add_text_frame(title_box, "THE FUTURE OF PRODUCT MARKETING", 14, WHITE, True)

# Blue gradient section (simulated with solid blue)
hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), Inches(13.333), Inches(4.5))
set_shape_fill(hero, SAP_BLUE)
hero.line.fill.background()

# Chapter tag
tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(1.5), Inches(0.4))
set_shape_fill(tag_box, SAP_ORANGE)
tag_box.line.fill.background()
tag_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.82), Inches(1.5), Inches(0.4))
tf = add_text_frame(tag_text, "CHAPTER 1", 10, WHITE, True, PP_ALIGN.CENTER)

# Main title
main_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(10), Inches(1.5))
tf = add_text_frame(main_title, "Your Buyer Just\nHired an Agent", 44, WHITE, True)

# Subtitle
subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(10), Inches(0.8))
add_text_frame(subtitle, "Everything PMMs assumed about buyers, positioning, and\ngo-to-market is breaking. Here's what actually changes.", 16, WHITE)

# Footer
footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
add_text_frame(footer, "Chris O'Hara  ·  futureofpmm.com  ·  March 2026", 11, DARK_SLATE)

# ========== SLIDE 2: The Assumption That Just Broke ==========
slide = prs.slides.add_slide(blank_layout)

# Header bar
header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "THE ASSUMPTION THAT JUST BROKE", 20, WHITE, True)

# Main content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(1.2))
tf = add_text_frame(content_box, "For twenty years, product marketing rested on one unspoken axiom:", 16, DARK_SLATE)

# Quote box
quote_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.3), Inches(12.3), Inches(1.2))
set_shape_fill(quote_bg, LIGHT_GRAY)
quote_bg.line.fill.background()
quote_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11.9), Inches(1))
add_text_frame(quote_text, '"The buyer is a human being who will, at some point in the purchase\nprocess, encounter the content you have carefully built for them."', 18, DARK_SLATE, False, PP_ALIGN.CENTER)

# Subtext
subtext = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12), Inches(0.6))
add_text_frame(subtext, "This assumption underpins every messaging framework, content strategy, SEO investment, and nurture sequence in your stack.", 14, DARK_SLATE)

# Three boxes showing old funnel
box_labels = ["Buyer searches\non Google", "Buyer reads content\ncase studies, demos", "Buyer engages\nfills form, books call"]
for i, label in enumerate(box_labels):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8 + i*4.2), Inches(4.8), Inches(3.8), Inches(1.5))
    set_shape_fill(box, SAP_BLUE)
    box.line.fill.background()
    text_box = slide.shapes.add_textbox(Inches(0.9 + i*4.2), Inches(5), Inches(3.6), Inches(1.3))
    tf = add_text_frame(text_box, label, 14, WHITE, False, PP_ALIGN.CENTER)
    # Checkmark
    check = slide.shapes.add_textbox(Inches(1 + i*4.2), Inches(6.3), Inches(3.4), Inches(0.4))
    add_text_frame(check, "✓ PMM content reaches this stage", 10, GREEN, False, PP_ALIGN.CENTER)

# ========== SLIDE 3: Meet Sarah ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "MEET SARAH — AND HER AGENT", 20, WHITE, True)

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
add_text_frame(subtitle, "A real scenario playing out at enterprise companies right now.", 14, DARK_SLATE)

# Timeline boxes
timeline_data = [
    ("Monday 9am", "Agent Dispatched", "Sarah's AI procurement assistant gets the brief: evaluate analytics platforms", SAP_BLUE),
    ("Monday 9:04", "Autonomous Research", "Queries structured DBs, G2 rankings, schema-marked content, analyst reports", SAP_BLUE),
    ("Monday 9:11", "Shortlist Generated", "3-vendor shortlist created. Comparison matrix built. Recommendation drafted.", SAP_BLUE),
    ("Monday 10am", "Sarah Reviews", "Opens her laptop. Shortlist already waiting. Vendors not on it: invisible.", SAP_ORANGE),
]

for i, (time, title, desc, color) in enumerate(timeline_data):
    x = Inches(0.5 + i*3.2)
    # Time label
    time_box = slide.shapes.add_textbox(x, Inches(2), Inches(3), Inches(0.4))
    add_text_frame(time_box, time, 12, DARK_SLATE, True, PP_ALIGN.CENTER)
    # Main box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.5), Inches(3), Inches(2.5))
    set_shape_fill(box, color)
    box.line.fill.background()
    # Title in box
    t_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.7), Inches(2.8), Inches(0.5))
    add_text_frame(t_box, title, 14, WHITE, True, PP_ALIGN.CENTER)
    # Description in box
    d_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.2), Inches(2.8), Inches(1.5))
    add_text_frame(d_box, desc, 11, WHITE, False, PP_ALIGN.CENTER)

# Warning callout
warning = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1))
set_shape_fill(warning, LIGHT_GRAY)
warning.line.fill.background()
warning_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.7))
add_text_frame(warning_text, "✗  Your case study, demo video, and product page were never seen.\n    The shortlist was set 2 minutes into Monday morning.", 14, DARK_SLATE, False, PP_ALIGN.LEFT)

# ========== SLIDE 4: The Visibility Gap ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "THE VISIBILITY GAP", 20, WHITE, True)

# Add image if exists
img_path = "/root/.openclaw/workspace/futureofpmm/textbook/ch1-visibility-gap.png"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), height=Inches(5))

caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
add_text_frame(caption, "The assets agents surface are already within PMM's scope — they just haven't been treated as strategic outputs.", 12, DARK_SLATE, False, PP_ALIGN.CENTER)

# ========== SLIDE 5: The Six Percent Problem ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "THE SIX PERCENT PROBLEM", 20, WHITE, True)

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
add_text_frame(subtitle, "After 35 editions and 180 articles synthesized — the signal is unmistakable.", 14, DARK_SLATE)

# Big number
big_num = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(3), Inches(2))
add_text_frame(big_num, "6%", 72, SAP_ORANGE, True)

label = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(3), Inches(0.5))
add_text_frame(label, "of B2B marketing\nteams are ready", 16, DARK_SLATE)

# Breakdown boxes
breakdown = [
    ("40%", "are unaware the transition is happening — still optimizing for Google", LIGHT_GRAY),
    ("35%", "are experimenting with tools but have no underlying strategy shift", LIGHT_GRAY),
    ("6%", "have made all three moves — and are pulling away from the field", SAP_ORANGE),
]

for i, (pct, desc, color) in enumerate(breakdown):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2 + i*1.5), Inches(9), Inches(1.3))
    set_shape_fill(box, color)
    box.line.fill.background()
    pct_text = slide.shapes.add_textbox(Inches(4.2), Inches(2.2 + i*1.5), Inches(1.2), Inches(1))
    add_text_frame(pct_text, pct, 24, DARK_SLATE if color == LIGHT_GRAY else WHITE, True)
    desc_text = slide.shapes.add_textbox(Inches(5.5), Inches(2.3 + i*1.5), Inches(7.3), Inches(1))
    add_text_frame(desc_text, desc, 14, DARK_SLATE if color == LIGHT_GRAY else WHITE)

# ========== SLIDE 6: Three Moves Overview ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "THE PMM RESPONSE: THREE MOVES", 20, WHITE, True)

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
add_text_frame(subtitle, "Winning PMMs operate all three layers simultaneously. Most organizations stall at Move 1, or skip directly to Move 3.", 14, DARK_SLATE)

# Three move boxes
moves = [
    ("MOVE 1", "Build the GEO Layer", "Generative Engine Optimization — structured data, schema markup, citation authority", SAP_BLUE),
    ("MOVE 2", "Build the Data Layer", "G2 profiles, structured pricing, analyst relations as data infrastructure", LIGHT_BLUE),
    ("MOVE 3", "Own the Narrative Layer", "Category framing, customer outcome narratives — the irreducibly human", SAP_ORANGE),
]

for i, (num, title, desc, color) in enumerate(moves):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5 + i*4.2), Inches(2.5), Inches(4), Inches(3.5))
    set_shape_fill(box, color)
    box.line.fill.background()
    
    num_text = slide.shapes.add_textbox(Inches(0.7 + i*4.2), Inches(2.7), Inches(3.6), Inches(0.5))
    add_text_frame(num_text, num, 12, WHITE, True, PP_ALIGN.CENTER)
    
    title_text = slide.shapes.add_textbox(Inches(0.7 + i*4.2), Inches(3.3), Inches(3.6), Inches(0.8))
    add_text_frame(title_text, title, 18, WHITE, True, PP_ALIGN.CENTER)
    
    desc_text = slide.shapes.add_textbox(Inches(0.7 + i*4.2), Inches(4.2), Inches(3.6), Inches(1.5))
    add_text_frame(desc_text, desc, 12, WHITE, False, PP_ALIGN.CENTER)

# ========== SLIDE 7: Move 1 - GEO Layer ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, SAP_BLUE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "MOVE 01  ·  BUILD THE GEO LAYER", 20, WHITE, True)

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
add_text_frame(subtitle, "Generative Engine Optimization is not SEO with a new name. The signals are fundamentally different.", 14, DARK_SLATE)

# Comparison table header
table_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(12.3), Inches(0.6))
set_shape_fill(table_header, DARK_SLATE)
table_header.line.fill.background()

cols = ["Signal", "SEO (Google)", "GEO (AI Agents)"]
for i, col in enumerate(cols):
    col_text = slide.shapes.add_textbox(Inches(0.6 + i*4.1), Inches(2.1), Inches(4), Inches(0.5))
    add_text_frame(col_text, col, 12, WHITE, True)

# Table rows
rows = [
    ("Optimizes for", "PageRank / inbound links", "Schema markup / citation authority"),
    ("Key content", "Keywords & density", "Structured FAQ + knowledge panels"),
    ("Review signals", "None directly", "G2 / Gartner Peer Insights (heavily weighted)"),
    ("Pricing", "Not a ranking factor", "Structured pricing = agent-processable"),
    ("Technical docs", "Helps crawlability", "Core agent evaluation signal"),
]

for i, (signal, seo, geo) in enumerate(rows):
    y = Inches(2.7 + i*0.6)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(0.55))
    set_shape_fill(row_bg, bg_color)
    row_bg.line.fill.background()
    
    for j, text in enumerate([signal, seo, geo]):
        cell = slide.shapes.add_textbox(Inches(0.6 + j*4.1), y + Inches(0.1), Inches(4), Inches(0.4))
        add_text_frame(cell, text, 11, DARK_SLATE, j==0)

# Quick win box
qw_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2))
set_shape_fill(qw_box, GREEN)
qw_box.line.fill.background()
qw_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(0.4))
add_text_frame(qw_title, "Quick Win · This Week", 12, WHITE, True)
qw_desc = slide.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(12), Inches(0.5))
add_text_frame(qw_desc, "Add structured FAQ markup (FAQPage schema) to your top 5 category pages. Single highest-ROI GEO action for most B2B PMM teams.", 12, WHITE)

# ========== SLIDE 8: Moves 2 & 3 ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "MOVE 02 & 03  ·  DATA LAYER + NARRATIVE LAYER", 20, WHITE, True)

# Move 2 box
m2_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(4.5))
set_shape_fill(m2_box, LIGHT_BLUE)
m2_box.line.fill.background()

m2_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5))
add_text_frame(m2_title, "02  Build the Data Layer", 16, WHITE, True)

m2_subtitle = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(5.6), Inches(0.4))
add_text_frame(m2_subtitle, "Structured Competitive Moat", 12, WHITE)

m2_bullets = [
    "Treat your G2 profile as a first-class PMM output",
    "Structure your pricing page with schema markup",
    "Analyst relations is now a data infrastructure project",
    "Competitive comparison pages with machine-readable markup",
]
for i, bullet in enumerate(m2_bullets):
    b = slide.shapes.add_textbox(Inches(0.7), Inches(2.6 + i*0.6), Inches(5.6), Inches(0.5))
    add_text_frame(b, "• " + bullet, 11, WHITE)

m2_time = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(5.6), Inches(0.4))
add_text_frame(m2_time, "⏱  30–60 day priority", 10, GOLD, True)

# Move 3 box
m3_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.3), Inches(6), Inches(4.5))
set_shape_fill(m3_box, SAP_ORANGE)
m3_box.line.fill.background()

m3_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.6), Inches(0.5))
add_text_frame(m3_title, "03  Own the Narrative Layer", 16, WHITE, True)

m3_subtitle = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.6), Inches(0.4))
add_text_frame(m3_subtitle, "Human-to-Human Irreducibles", 12, WHITE)

m3_bullets = [
    "This is what AI cannot replace — and the noise floor rising makes it more valuable",
    "Category framing: make your position feel inevitable, not chosen",
    "Customer outcome narratives that only you have the relationships to make real",
    "Moves 1 & 2 get you in the room. Move 3 is how you win.",
]
for i, bullet in enumerate(m3_bullets):
    b = slide.shapes.add_textbox(Inches(7), Inches(2.6 + i*0.65), Inches(5.6), Inches(0.6))
    add_text_frame(b, "• " + bullet, 11, WHITE)

m3_time = slide.shapes.add_textbox(Inches(7), Inches(5.2), Inches(5.6), Inches(0.4))
add_text_frame(m3_time, "⏱  Ongoing investment", 10, GOLD, True)

# ========== SLIDE 9: What To Do This Week ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "WHAT TO DO THIS WEEK", 20, WHITE, True)

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
add_text_frame(subtitle, "The gap is compounding. But it closes with specific, sequenced moves — not a multi-quarter strategy.", 14, DARK_SLATE)

# Action boxes
actions = [
    ("This week", "Run the Agent Audit", "Ask an AI assistant to research your category and find top vendors. Note where you appear, what it says, and what it cites.", SAP_BLUE),
    ("Next 30 days", "Implement Schema Markup", "FAQPage and SoftwareApplication schemas on your top 5 pages. Highest-ROI GEO action.", SAP_BLUE),
    ("Next 60 days", "Audit G2 + Gartner Presence", "Run a review generation campaign. Ensure your profiles surface the right features.", LIGHT_BLUE),
    ("Next 90 days", "Build Structured Competitive Content", "One genuine comparison page with machine-readable markup.", SAP_ORANGE),
]

for i, (timing, title, desc, color) in enumerate(actions):
    x = Inches(0.5 + i*3.2)
    
    time_box = slide.shapes.add_textbox(x, Inches(2), Inches(3), Inches(0.4))
    add_text_frame(time_box, timing, 11, DARK_SLATE, True, PP_ALIGN.CENTER)
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.5), Inches(3), Inches(3.5))
    set_shape_fill(box, color)
    box.line.fill.background()
    
    t = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.7), Inches(2.8), Inches(0.7))
    add_text_frame(t, title, 14, WHITE, True, PP_ALIGN.CENTER)
    
    d = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.4), Inches(2.8), Inches(2.4))
    add_text_frame(d, desc, 11, WHITE, False, PP_ALIGN.CENTER)

# Footer
footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), Inches(13.333), Inches(0.7))
set_shape_fill(footer_bar, LIGHT_GRAY)
footer_bar.line.fill.background()
footer_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12), Inches(0.4))
add_text_frame(footer_text, "futureofpmm.com  ·  Chris O'Hara  ·  The Future of Product Marketing", 10, DARK_SLATE, False, PP_ALIGN.CENTER)

# ========== SLIDE 10: CMO Perspective ==========
slide = prs.slides.add_slide(blank_layout)

header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_shape_fill(header, DARK_SLATE)
header.line.fill.background()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text_frame(title_box, "CMO PERSPECTIVE", 20, WHITE, True)

# Main quote
quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2))
set_shape_fill(quote_box, LIGHT_GRAY)
quote_box.line.fill.background()
quote = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.8))
add_text_frame(quote, "The conversation about AI inside marketing departments is almost entirely wrong.\nIt's focused on efficiency — how to produce the same outputs with fewer people.\n\nThat's a CFO question, not a CMO question.\n\nThe CMO question is: how do we produce fundamentally different outputs?", 14, DARK_SLATE, False, PP_ALIGN.CENTER)

# Key takeaways
takeaways_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12), Inches(0.5))
add_text_frame(takeaways_title, "KEY TAKEAWAYS", 14, DARK_SLATE, True)

takeaways = [
    "The agentic shift is a structural transformation, not a toolkit upgrade",
    "Buyer behavior has already changed — agents compress research from weeks to hours",
    "The \"10x\" thesis isn't about speed — it's about the widening output gap",
    "PMMs who internalize the shift will pull away from those who don't",
]

for i, t in enumerate(takeaways):
    bullet_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.1 + i*0.6), Inches(12.3), Inches(0.5))
    set_shape_fill(bullet_box, SAP_BLUE if i < 2 else SAP_ORANGE)
    bullet_box.line.fill.background()
    bullet_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.2 + i*0.6), Inches(11.9), Inches(0.4))
    add_text_frame(bullet_text, t, 12, WHITE)

# Save
output_path = "/root/.openclaw/workspace/futureofpmm/textbook/chapter-01-sap.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
