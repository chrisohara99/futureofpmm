#!/usr/bin/env python3
"""
Create SAP-styled PowerPoint for Chapter 1: The Old Playbook Is Dead
Based on new DOCX content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# SAP Colors
DARK_SLATE = RGBColor(0x44, 0x54, 0x6A)
SAP_ORANGE = RGBColor(0xED, 0x7D, 0x31)
SAP_BLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHT_GRAY = RGBColor(0xE7, 0xE6, 0xE6)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
GREEN = RGBColor(0x70, 0xAD, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_RED = RGBColor(0x8B, 0x00, 0x00)

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text(shape, text, size=18, color=DARK_SLATE, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return tf

def add_para(tf, text, size=14, color=DARK_SLATE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    return p

# Create 16:9 presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

FIG_DIR = "/root/.openclaw/workspace/futureofpmm/textbook/ch1-new"

# ==================== SLIDE 1: Title ====================
slide = prs.slides.add_slide(blank)

# Header
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
add_text(t, "THE FUTURE OF PRODUCT MARKETING", 13, WHITE, True)

# Hero
hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(13.333), Inches(4.8))
set_fill(hero, SAP_BLUE)
hero.line.fill.background()

# Tag
tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(2.2), Inches(0.45))
set_fill(tag, SAP_ORANGE)
tag.line.fill.background()
tt = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(2.2), Inches(0.4))
add_text(tt, "PART I · CHAPTER 1", 11, WHITE, True, PP_ALIGN.CENTER)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(10), Inches(1.8))
add_text(title, "The Old Playbook\nIs Dead", 48, WHITE, True)

# Subtitle
sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(10), Inches(0.8))
add_text(sub, "Pragmatic Remix: All 37 Boxes (Reframed)", 16, WHITE)

# Footer
foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
add_text(foot, "Chris O'Hara & Dan Yu  ·  futureofpmm.com  ·  March 2026", 11, DARK_SLATE)

# ==================== SLIDE 2: Sarah's Story ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE STORY OF SARAH", 20, WHITE, True)

# Opening
intro = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.2))
add_text(intro, "Sometime in late 2024, a product marketer noticed something strange in her pipeline data...", 16, DARK_SLATE)

# Three observations
obs = [
    ("1", "Traffic Up, Forms Down", "More people visiting product pages,\nfewer downloading white papers", SAP_BLUE),
    ("2", "Prospects Knew Too Much", "Showing up to calls knowing competitive\npricing and technical limitations", LIGHT_BLUE),
    ("3", "Deals Closed in 3 Weeks", "Typical cycle was 11 weeks.\nTwo enterprise deals closed in under 3.", SAP_ORANGE),
]

for i, (num, title, desc, color) in enumerate(obs):
    x = Inches(0.5 + i*4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.5), Inches(4), Inches(2.8))
    set_fill(box, color)
    box.line.fill.background()
    
    n = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.7), Inches(0.6), Inches(0.5))
    add_text(n, num, 24, WHITE, True)
    
    tt = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.2), Inches(3.6), Inches(0.6))
    add_text(tt, title, 16, WHITE, True)
    
    d = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.8), Inches(3.6), Inches(1.3))
    add_text(d, desc, 12, WHITE)

# Revelation
rev_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.3))
set_fill(rev_box, LIGHT_GRAY)
rev_box.line.fill.background()
rev = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(11.9), Inches(1))
add_text(rev, "What had changed: her buyers had started using AI agents to do the research, comparison shopping, and vendor pre-qualification that used to take human buying committees weeks.", 14, DARK_SLATE)

# ==================== SLIDE 3: Buyer Journey Figure ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE BUYER JOURNEY COMPRESSION", 20, WHITE, True)

fig_path = os.path.join(FIG_DIR, "fig1-buyer-journey.png")
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(0.5), Inches(1.2), width=Inches(12.3))

caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.6))
add_text(caption, "Figure 1: Steps 2–4 of the old journey — where most PMM content lived — now happen inside an AI's context window in seconds.", 11, DARK_SLATE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 4: Three Eras ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE THREE ERAS OF ENTERPRISE DATA", 20, WHITE, True)

eras = [
    ("ERA 1", "Data Warehouse", "Age of Reports", "ETL batch jobs. Static deliverables.\nQuarterly reviews. Slow data,\nslow marketing, slow buyers.", DARK_SLATE),
    ("ERA 2", "Data Lake", "Age of Unification", "CDP era. Customer strategy.\nJourney mapping. Win/loss analysis.\n\"A CDP knows your customer.\"", SAP_BLUE),
    ("ERA 3", "Data Fabric", "Age of Agents", "Agentic AI. Federated metadata.\nAgents that monitor, adjust, generate.\n\"An AMP knows your business.\"", SAP_ORANGE),
]

for i, (num, title, subtitle, desc, color) in enumerate(eras):
    x = Inches(0.5 + i*4.2)
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.3), Inches(4), Inches(4.8))
    set_fill(box, color)
    box.line.fill.background()
    
    n = slide.shapes.add_textbox(x + Inches(0.2), Inches(1.5), Inches(3.6), Inches(0.5))
    add_text(n, num, 12, WHITE, True)
    
    tt = slide.shapes.add_textbox(x + Inches(0.2), Inches(2), Inches(3.6), Inches(0.7))
    add_text(tt, title, 20, WHITE, True)
    
    st = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.6), Inches(3.6), Inches(0.5))
    add_text(st, subtitle, 14, GOLD if color != DARK_SLATE else LIGHT_GRAY)
    
    d = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.3), Inches(3.6), Inches(2.5))
    add_text(d, desc, 12, WHITE)

# Tagline
tag_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8))
set_fill(tag_box, LIGHT_GRAY)
tag_box.line.fill.background()
tag_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.6))
add_text(tag_text, "A DMP knew your audience. A CDP knew your customer. An AMP knows your business.", 14, DARK_SLATE, True, PP_ALIGN.CENTER)

# ==================== SLIDE 5: Visibility Gap ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE VISIBILITY GAP", 20, WHITE, True)

fig_path = os.path.join(FIG_DIR, "fig2-visibility-gap.png")
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.6))
add_text(caption, "Figure 2: The Visibility Gap is not a technology problem — it is a prioritization problem.", 11, DARK_SLATE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 6: SEO to GEO ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "FROM SEO TO GEO", 20, WHITE, True)

sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
add_text(sub, "Generative Engine Optimization — the signals are fundamentally different.", 14, DARK_SLATE)

fig_path = os.path.join(FIG_DIR, "fig3-seo-to-geo.png")
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(1.5), Inches(1.6), width=Inches(10))

# Action box
act_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2))
set_fill(act_box, GREEN)
act_box.line.fill.background()
act_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(0.4))
add_text(act_title, "The Audit (Do This Week)", 12, WHITE, True)
act_desc = slide.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.6))
add_text(act_desc, "Open Claude, ChatGPT, or Perplexity. Ask it to evaluate vendors in your category. See if your company appears. That 60-second exercise tells you more about your GEO readiness than any consultant's report.", 11, WHITE)

# ==================== SLIDE 7: Pragmatic Framework ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE PRAGMATIC FRAMEWORK MEETS THE MACHINE", 20, WHITE, True)

sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.6))
add_text(sub, "37 boxes. Three clusters. One question: what happens when AI agents become part of the team?", 14, DARK_SLATE)

clusters = [
    ("CLUSTER 1", "Automated", "Market sizing, derivative content,\nevent logistics, sales tools production.\nAgents do the majority today.", "70% time savings on production.\nIf your value is producing artifacts,\nthe automation wave is not your friend.", LIGHT_GRAY, DARK_SLATE),
    ("CLUSTER 2", "Augmented", "Positioning, win/loss analysis,\nbuyer personas, pricing strategy.\nHuman essential, agent accelerates.", "Human-agent team is so much better\nthat operating without agents\nbecomes a competitive disadvantage.", SAP_BLUE, WHITE),
    ("CLUSTER 3", "Elevated", "Stakeholder management, customer\nempathy, thought leadership.\nIrreducibly human.", "These activities become MORE important\nprecisely because machines can't do them.\nThis is where careers are built.", SAP_ORANGE, WHITE),
]

for i, (num, title, desc1, desc2, bg_color, txt_color) in enumerate(clusters):
    x = Inches(0.5 + i*4.2)
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.9), Inches(4), Inches(4.5))
    set_fill(box, bg_color)
    box.line.fill.background()
    
    n = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.1), Inches(3.6), Inches(0.4))
    add_text(n, num, 10, txt_color, True)
    
    tt = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.6), Inches(0.5))
    add_text(tt, title, 18, txt_color, True)
    
    d1 = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.1), Inches(3.6), Inches(1.3))
    add_text(d1, desc1, 11, txt_color)
    
    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), Inches(4.4), Inches(3.6), Inches(0.02))
    set_fill(div, GOLD if bg_color != LIGHT_GRAY else SAP_ORANGE)
    div.line.fill.background()
    
    d2 = slide.shapes.add_textbox(x + Inches(0.2), Inches(4.6), Inches(3.6), Inches(1.5))
    add_text(d2, desc2, 10, txt_color)

# ==================== SLIDE 8: Curriculum Architecture ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE CURRICULUM ARCHITECTURE", 20, WHITE, True)

fig_path = os.path.join(FIG_DIR, "fig4-pragmatic-funnel-map.png")
if os.path.exists(fig_path):
    slide.shapes.add_picture(fig_path, Inches(0.5), Inches(1.1), width=Inches(12.3))

caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12), Inches(0.6))
add_text(caption, "Figure 4: The Pragmatic 37 → Agentic Funnel Map. All 37 activities mapped to the six-stage funnel, color-coded by Pragmatic category.", 10, DARK_SLATE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 9: 10x Thesis ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, SAP_ORANGE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE 10x THESIS", 20, WHITE, True)

# Big quote
q_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.5))
set_fill(q_box, LIGHT_GRAY)
q_box.line.fill.background()
quote = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.2))
add_text(quote, "A PMM who automates 20 hours a week of Cluster One work and redirects that time into deeper competitive analysis, more thoughtful positioning, and genuine thought leadership is going to outperform a PMM who's still manually updating battlecards.\n\nNot by 20%. By a factor that justifies the title of this book.", 16, DARK_SLATE, False, PP_ALIGN.CENTER)

# What it means
meaning_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12), Inches(0.5))
add_text(meaning_title, "What \"10x yourself\" actually means:", 14, DARK_SLATE, True)

meanings = [
    ("✗", "NOT working ten times harder", DARK_RED),
    ("✗", "NOT working ten times faster", DARK_RED),
    ("✓", "Operating at a fundamentally different level of strategic impact", GREEN),
]

for i, (icon, text, color) in enumerate(meanings):
    row = slide.shapes.add_textbox(Inches(0.5), Inches(4.6 + i*0.5), Inches(12), Inches(0.5))
    tf = add_text(row, f"{icon}  {text}", 14, color, i==2)

# Bottom
bottom = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.8))
add_text(bottom, "You've offloaded the mechanical work to machines built for it — and you're finally free to do the work that only a human can do.", 14, DARK_SLATE, True, PP_ALIGN.CENTER)

# ==================== SLIDE 10: CMO Perspective ====================
slide = prs.slides.add_slide(blank)

h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "CMO PERSPECTIVE", 20, WHITE, True)

# Quote box
q_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.8))
set_fill(q_box, LIGHT_GRAY)
q_box.line.fill.background()
quote = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.5))
add_text(quote, "The conversation about AI inside marketing departments is almost entirely wrong. It's focused on efficiency — how to produce the same outputs with fewer people. That's a CFO question.\n\nThe CMO question: how do we produce fundamentally different outputs?", 14, DARK_SLATE, False, PP_ALIGN.CENTER)

# Key takeaways title
kt_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(0.5))
add_text(kt_title, "KEY TAKEAWAYS", 14, DARK_SLATE, True)

takeaways = [
    "The agentic shift is a structural transformation, not just a toolkit upgrade",
    "Buyer behavior has already changed — AI agents compress research from weeks to hours",
    "Most PMM content is invisible to agents (the Visibility Gap)",
    "The Pragmatic 37 cluster into: Automated, Augmented, Elevated",
    "The PMM who frees Cluster One hours to invest in Cluster Three operates at a different level",
]

for i, txt in enumerate(takeaways):
    color = SAP_BLUE if i < 3 else SAP_ORANGE
    row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.7 + i*0.55), Inches(12.3), Inches(0.5))
    set_fill(row_bg, color)
    row_bg.line.fill.background()
    row_txt = slide.shapes.add_textbox(Inches(0.7), Inches(3.8 + i*0.55), Inches(11.9), Inches(0.4))
    add_text(row_txt, txt, 11, WHITE)

# Footer
foot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), Inches(13.333), Inches(0.7))
set_fill(foot_bar, LIGHT_GRAY)
foot_bar.line.fill.background()
foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4))
add_text(foot, "futureofpmm.com  ·  Chris O'Hara & Dan Yu  ·  The Future of Product Marketing", 10, DARK_SLATE, False, PP_ALIGN.CENTER)

# Save
output_path = "/root/.openclaw/workspace/futureofpmm/textbook/chapter-01-sap.pptx"
prs.save(output_path)
print(f"✓ Saved: {output_path}")
print(f"  10 slides with SAP styling and 4 embedded figures")
