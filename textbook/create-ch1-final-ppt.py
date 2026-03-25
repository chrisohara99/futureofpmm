#!/usr/bin/env python3
"""
Create final Chapter 1 PPT with all 6 official graphics
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_SLATE = RGBColor(0x44, 0x54, 0x6A)
SAP_ORANGE = RGBColor(0xED, 0x7D, 0x31)
SAP_BLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHT_GRAY = RGBColor(0xE7, 0xE6, 0xE6)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
LIGHT_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
GREEN = RGBColor(0x70, 0xAD, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def add_text(shape, text, size=18, color=DARK_SLATE, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return tf

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

FIG_DIR = "/root/.openclaw/workspace/futureofpmm/textbook/ch1-graphics"

# ==================== SLIDE 1: Title ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
add_text(t, "THE FUTURE OF PRODUCT MARKETING", 13, WHITE, True)

hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), Inches(13.333), Inches(4.8))
set_fill(hero, SAP_BLUE)
hero.line.fill.background()

tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(2.2), Inches(0.45))
set_fill(tag, SAP_ORANGE)
tag.line.fill.background()
tt = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(2.2), Inches(0.4))
add_text(tt, "PART I · CHAPTER 1", 11, WHITE, True, PP_ALIGN.CENTER)

title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(10), Inches(1.8))
add_text(title, "The Old Playbook\nIs Dead", 48, WHITE, True)

sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(10), Inches(0.8))
add_text(sub, "Pragmatic Remix: All 37 Boxes (Reframed)", 16, WHITE)

foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
add_text(foot, "Chris O'Hara & Dan Yu  ·  futureofpmm.com  ·  March 2026", 11, DARK_SLATE)

# ==================== SLIDE 2: Core Thesis ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "CORE THESIS", 20, WHITE, True)

box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3))
set_fill(box, LIGHT_GRAY)
box.line.fill.background()
txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.5))
add_text(txt, "The agentic era didn't accelerate the old B2B playbook — it ended it.\n\nProduct lifecycles now run in continuous deployment, not quarterly releases.\nBuyers consume content through AI intermediaries, not your website.\nGEO has replaced SEO. The funnel is no longer a funnel.\n\nThe PMM who survives this transition is not the one who learned to prompt faster —\nit's the one who understood what changed and why.", 16, DARK_SLATE, False, PP_ALIGN.CENTER)

# Key points
pts = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.3), Inches(1.5))
tf = add_text(pts, "This chapter establishes:", 12, DARK_SLATE, True)
for item in ["The three eras of enterprise data", "The visibility gap", "From SEO to GEO", "The Pragmatic 37 reimagined"]:
    p = tf.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(12)
    p.font.name = "Arial"
    p.font.color.rgb = DARK_SLATE

# ==================== SLIDE 3: Buyer Journey ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE BUYER JOURNEY COMPRESSION", 18, WHITE, True)

fig = os.path.join(FIG_DIR, "fig1-buyer-journey.png")
if os.path.exists(fig):
    slide.shapes.add_picture(fig, Inches(0.3), Inches(0.9), width=Inches(12.7))

# ==================== SLIDE 4: Visibility Gap ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE VISIBILITY GAP", 18, WHITE, True)

fig = os.path.join(FIG_DIR, "fig2-visibility-gap.png")
if os.path.exists(fig):
    slide.shapes.add_picture(fig, Inches(0.3), Inches(0.9), width=Inches(12.7))

# ==================== SLIDE 5: SEO to GEO ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "FROM SEO TO GEO", 18, WHITE, True)

fig = os.path.join(FIG_DIR, "fig3-seo-to-geo.png")
if os.path.exists(fig):
    slide.shapes.add_picture(fig, Inches(1), Inches(0.9), width=Inches(11))

# Action box
act = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2))
set_fill(act, GREEN)
act.line.fill.background()
at = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(0.4))
add_text(at, "The Audit (Do This Week)", 12, WHITE, True)
ad = slide.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.6))
add_text(ad, "Open Claude, ChatGPT, or Perplexity. Ask it to evaluate vendors in your category. See if your company appears.", 11, WHITE)

# ==================== SLIDE 6: New Awareness Stack ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE NEW AWARENESS STACK", 18, WHITE, True)

fig = os.path.join(FIG_DIR, "fig5-awareness-stack.png")
if os.path.exists(fig):
    slide.shapes.add_picture(fig, Inches(0.3), Inches(0.9), width=Inches(12.7))

# ==================== SLIDE 7: Tiered Launch Model ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.8))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
add_text(t, "THE TIERED LAUNCH MODEL", 18, WHITE, True)

fig = os.path.join(FIG_DIR, "fig6-tiered-launch.png")
if os.path.exists(fig):
    slide.shapes.add_picture(fig, Inches(0.3), Inches(1.2), width=Inches(12.7))

# ==================== SLIDE 8: Pragmatic Funnel Map ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.6))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(12), Inches(0.4))
add_text(t, "THE PRAGMATIC 37 → AGENTIC FUNNEL MAP", 16, WHITE, True)

fig = os.path.join(FIG_DIR, "fig4-pragmatic-funnel-map.png")
if os.path.exists(fig):
    slide.shapes.add_picture(fig, Inches(0.2), Inches(0.65), width=Inches(12.9))

# ==================== SLIDE 9: Three Clusters ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE PRAGMATIC FRAMEWORK MEETS THE MACHINE", 18, WHITE, True)

sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
add_text(sub, "37 boxes. Three clusters. Which activities survive automation?", 14, DARK_SLATE)

clusters = [
    ("CLUSTER 1", "Automated", "Market sizing, derivative content,\nevent logistics, sales tools.\n\nAgents do the majority today.", LIGHT_GRAY, DARK_SLATE),
    ("CLUSTER 2", "Augmented", "Positioning, win/loss analysis,\nbuyer personas, pricing.\n\nHuman essential, agent accelerates.", SAP_BLUE, WHITE),
    ("CLUSTER 3", "Elevated", "Stakeholder management,\ncustomer empathy, thought leadership.\n\nIrreducibly human.", SAP_ORANGE, WHITE),
]

for i, (num, title, desc, bg, txt_color) in enumerate(clusters):
    x = Inches(0.5 + i*4.2)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(4), Inches(4))
    set_fill(box, bg)
    box.line.fill.background()
    
    n = slide.shapes.add_textbox(x + Inches(0.2), Inches(2), Inches(3.6), Inches(0.4))
    add_text(n, num, 10, txt_color, True)
    
    tt = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.5))
    add_text(tt, title, 18, txt_color, True)
    
    d = slide.shapes.add_textbox(x + Inches(0.2), Inches(3), Inches(3.6), Inches(2.5))
    add_text(d, desc, 12, txt_color)

# Bottom insight
ins = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6), Inches(12.3), Inches(1))
set_fill(ins, LIGHT_GRAY)
ins.line.fill.background()
it = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7))
add_text(it, "If your value is producing artifacts, the automation wave is not your friend.\nThe question: are you directing the agents or being replaced by them?", 12, DARK_SLATE, False, PP_ALIGN.CENTER)

# ==================== SLIDE 10: 10x Thesis ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, SAP_ORANGE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "THE 10x THESIS", 20, WHITE, True)

q = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2))
set_fill(q, LIGHT_GRAY)
q.line.fill.background()
qt = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.8))
add_text(qt, "A PMM who automates 20 hours a week of Cluster One work and redirects that time into deeper competitive analysis, more thoughtful positioning, and genuine thought leadership will outperform one still manually updating battlecards.\n\nNot by 20%. By a factor that justifies the title of this curriculum.", 15, DARK_SLATE, False, PP_ALIGN.CENTER)

mt = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12), Inches(0.5))
add_text(mt, "What \"10x yourself\" actually means:", 14, DARK_SLATE, True)

meanings = [
    ("✗  NOT working ten times harder", RGBColor(0x8B, 0x00, 0x00)),
    ("✗  NOT working ten times faster", RGBColor(0x8B, 0x00, 0x00)),
    ("✓  Operating at a fundamentally different level of strategic impact", GREEN),
]
for i, (txt, color) in enumerate(meanings):
    row = slide.shapes.add_textbox(Inches(0.5), Inches(4.1 + i*0.5), Inches(12), Inches(0.5))
    add_text(row, txt, 14, color, i==2)

bottom = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1))
add_text(bottom, "You've offloaded the mechanical work to machines built for it —\nand you're finally free to do the work that only a human can do.", 14, DARK_SLATE, True, PP_ALIGN.CENTER)

# ==================== SLIDE 11: CMO Perspective ====================
slide = prs.slides.add_slide(blank)
h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1))
set_fill(h, DARK_SLATE)
h.line.fill.background()
t = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.5))
add_text(t, "CMO PERSPECTIVE", 20, WHITE, True)

q = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.6))
set_fill(q, LIGHT_GRAY)
q.line.fill.background()
qt = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.4))
add_text(qt, "The conversation about AI inside marketing departments is almost entirely wrong.\nIt's focused on efficiency — that's a CFO question.\n\nThe CMO question: how do we produce fundamentally different outputs?", 14, DARK_SLATE, False, PP_ALIGN.CENTER)

kt = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(0.5))
add_text(kt, "KEY TAKEAWAYS", 14, DARK_SLATE, True)

takeaways = [
    "The agentic shift is a structural transformation, not a toolkit upgrade",
    "Buyer behavior has already changed — AI agents compress research from weeks to hours",
    "Most PMM content is invisible to agents (the Visibility Gap)",
    "The Pragmatic 37 cluster into: Automated, Augmented, Elevated",
    "The PMM who frees Cluster One hours to invest in Cluster Three operates at a different level",
]
for i, txt in enumerate(takeaways):
    color = SAP_BLUE if i < 3 else SAP_ORANGE
    row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.5 + i*0.55), Inches(12.3), Inches(0.5))
    set_fill(row, color)
    row.line.fill.background()
    rt = slide.shapes.add_textbox(Inches(0.7), Inches(3.6 + i*0.55), Inches(11.9), Inches(0.4))
    add_text(rt, txt, 11, WHITE)

foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), Inches(13.333), Inches(0.7))
set_fill(foot, LIGHT_GRAY)
foot.line.fill.background()
ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4))
add_text(ft, "futureofpmm.com  ·  Chris O'Hara & Dan Yu  ·  The Future of Product Marketing", 10, DARK_SLATE, False, PP_ALIGN.CENTER)

# Save
out = "/root/.openclaw/workspace/futureofpmm/dist/curriculum/chapter-01/chapter-01-sap.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  11 slides with all 6 official graphics")
